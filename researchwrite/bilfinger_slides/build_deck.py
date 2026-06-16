#!/usr/bin/env python3
"""
Build the barrel-detection thesis slide deck (.pptx) from rendered figures.

Run make_figures.py first (writes assets/*.png), then this.
Output: researchwrite/bilfinger_slides/Bilfinger_barrel_detection_progress.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "Semi-Occluded_Object_Detection_Progress.pptx")
DATELINE = "June 2026"

THESIS_TITLE = "Detection and Manipulation of Semi-Occluded Objects"
THESIS_SUB = ("Barrel detection in 3D point clouds — benchmarking LiDAR-only "
              "and camera–LiDAR fusion methods")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

DARK = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x0B, 0x5C, 0x8A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def setp(p, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         space_after=6, level=0, italic=False):
    p.text = text
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"


def add_para(tf, *a, **k):
    p = tf.add_paragraph()
    setp(p, *a, **k)
    return p


def bullet(tf, text, size=15, color=DARK, level=0, bold=False, first=False,
           space_after=6, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    pre = "• " if level == 0 else "– "
    setp(p, pre + text, size=size, bold=bold, color=color, level=level,
         space_after=space_after, italic=italic)
    return p


def band(s, color=ACCENT, h=Inches(1.15)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def title_band(s, title, sub=None):
    band(s)
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.95))
    setp(tf.paragraphs[0], title, size=28, bold=True, color=WHITE)
    if sub:
        add_para(tf, sub, size=14, color=RGBColor(0xD6, 0xE6, 0xF0))


def footer(s, n):
    tf = box(s, Inches(0.4), Inches(7.05), Inches(12.0), Inches(0.35))
    setp(tf.paragraphs[0], THESIS_TITLE, size=9, color=GREY)
    tn = box(s, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35))
    setp(tn.paragraphs[0], str(n), size=10, color=GREY, align=PP_ALIGN.RIGHT)


def pic_fit(s, path, l, t, w, h):
    from PIL import Image
    iw, ih = Image.open(path).size
    boxr, imgr = w / h, iw / ih
    if imgr > boxr:
        nw, nh = w, int(w / imgr)
    else:
        nh, nw = h, int(h * imgr)
    nl = l + (w - nw) // 2
    nt = t + (h - nh) // 2
    s.shapes.add_picture(path, Emu(int(nl)), Emu(int(nt)),
                         width=Emu(int(nw)), height=Emu(int(nh)))


def dark_bg(s):
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK
    bg.line.fill.background(); bg.shadow.inherit = False


# =====================================================================
# 1. TITLE  (thesis name only)
# =====================================================================
s = slide()
dark_bg(s)
band(s, ACCENT, Inches(0.18))
tf = box(s, Inches(0.9), Inches(2.4), Inches(11.6), Inches(2.8))
setp(tf.paragraphs[0], THESIS_TITLE, size=36, bold=True, color=WHITE)
add_para(tf, THESIS_SUB, size=22, color=RGBColor(0x9F, 0xC5, 0xDD),
         space_after=18)
tf2 = box(s, Inches(0.9), Inches(5.5), Inches(11.6), Inches(1.4))
setp(tf2.paragraphs[0], "Master's thesis  ·  Bharath  ·  " + DATELINE,
     size=16, color=RGBColor(0xC9, 0xD6, 0xE0))

# =====================================================================
# 2. APPROACH  (the idea: standardised benchmark -> Isaac Sim)
# =====================================================================
s = slide()
title_band(s, "Approach: a method-agnostic benchmark, scaled to simulation")
tf = box(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5))
bullet(tf, "Standardised I/O contract: every detector consumes the same "
           "point cloud and emits a parametric cylinder "
           "(radius, axis, axis-point, extent) — evaluation is therefore "
           "method-agnostic and identical across methods.",
       size=15, first=True, bold=True)
bullet(tf, "Stage 1 — establish and validate the metric fit on controlled "
           "input:", size=15, bold=True, color=ACCENT, space_after=3)
bullet(tf, "synthetic clouds with analytic ground truth, plus real depth "
           "captures for sim-to-real transfer.", size=14, level=1)
bullet(tf, "Stage 2 — scale ground-truthed data generation with "
           "NVIDIA Isaac Sim:", size=15, bold=True, color=ACCENT, space_after=3)
bullet(tf, "per-barrel pose/size/count, instance & semantic segmentation, "
           "RTX LiDAR matched to the real sensor with a co-registered camera; "
           "scripted occlusion, burial depth and inter-barrel spacing.",
       size=14, level=1)
bullet(tf, "Stage 3 — extend the benchmark from LiDAR-only geometric to "
           "learned detectors and camera–LiDAR fusion, scored on the same "
           "metrics.", size=15, bold=True, color=ACCENT)
bullet(tf, "Metrics: precision / recall / F1, radius RMSE, axis-angle error, "
           "recall vs. occlusion fraction, runtime; annotation/training cost "
           "tracked so geometric and learned methods compare fairly.",
       size=14, color=GREY)
footer(s, 2)

def caption(s, text, l, t, w, color=GREY):
    tf = box(s, l, t, w, Inches(0.35))
    setp(tf.paragraphs[0], text, size=12, color=color, align=PP_ALIGN.CENTER)


# =====================================================================
# 3. SYNTHETIC validation  (first) — 3D render + 2D quantitative overlay
# =====================================================================
s = slide()
title_band(s, "Validation on synthetic clouds with analytic ground truth",
           "Generated cylinders with analytic radial normals — exact "
           "(radius, axis, centre) known")
pic_fit(s, os.path.join(ASSETS, "synth_3d_det_gt.png"),
        Inches(0.5), Inches(1.55), Inches(3.0), Inches(4.7))
caption(s, "3D: fitted barrel + synthetic points", Inches(0.3), Inches(6.3),
        Inches(3.4))
pic_fit(s, os.path.join(ASSETS, "synth_det_vs_gt.png"),
        Inches(3.7), Inches(1.55), Inches(4.8), Inches(4.7))
caption(s, "2D overlay: detected (red) vs ground truth (green)",
        Inches(3.7), Inches(6.3), Inches(4.8))
tf = box(s, Inches(8.8), Inches(1.7), Inches(4.2), Inches(5.0))
bullet(tf, "Red = detected cylinder, green = ground truth — near-exact "
           "overlap.", size=14, first=True, color=GREEN)
bullet(tf, "Set spans a one-sided ~120° arc (Xtion-like partial view) and an "
           "additive Gaussian-noise sweep (σ = n0.1–0.3).", size=14)
bullet(tf, "Top methods: P = R = F1 = 1.00; radius RMSE < 0.5 cm; "
           "axis-angle error < 0.5°.", size=14, color=GREEN)
bullet(tf, "Upper bound on fit accuracy where truth is exact, before "
           "real-sensor noise.", size=14)
add_para(tf, "SYNTHETIC DATA", size=14, bold=True, color=AMBER)
footer(s, 3)

# =====================================================================
# 4. REAL transfer — 3D raw vs detected
# =====================================================================
s = slide()
title_band(s, "Transfer to real depth captures",
           "Asus Xtion Pro, single barrel, ~180° of the surface visible "
           "(one-sided view)")
pic_fit(s, os.path.join(ASSETS, "real_3d_raw.png"),
        Inches(0.6), Inches(1.55), Inches(3.0), Inches(4.7))
caption(s, "raw depth cloud", Inches(0.6), Inches(6.3), Inches(3.0))
pic_fit(s, os.path.join(ASSETS, "real_3d_det.png"),
        Inches(4.0), Inches(1.55), Inches(3.0), Inches(4.7))
caption(s, "fitted barrel (red) + measured points", Inches(3.5), Inches(6.3),
        Inches(4.0))
tf = box(s, Inches(7.7), Inches(1.7), Inches(5.2), Inches(5.0))
bullet(tf, "Pipeline: depth → surface-normal estimation (Open3D) → "
           "cylinder fit.", size=15, first=True)
bullet(tf, "RANSAC fit: radius 4.63 cm vs. 4.25 cm physical "
           "(error 0.38 cm).", size=15, color=GREEN)
bullet(tf, "Axis-angle error 1.85° against the vertical reference.",
       size=15, color=GREEN)
bullet(tf, "Real-sensor radius error ~6–8× the synthetic figure — "
           "quantifies the sim-to-real gap.", size=15)
add_para(tf, "REAL DATA", size=14, bold=True, color=AMBER)
footer(s, 4)

# =====================================================================
# 5. FOUR METHODS
# =====================================================================
s = slide()
title_band(s, "Four LiDAR-only geometric detectors on identical input",
           "Same real barrel cloud — four independent algorithms")
pic_fit(s, os.path.join(ASSETS, "methods_grid_real.png"),
        Inches(0.4), Inches(1.5), Inches(12.5), Inches(3.6))
tf = box(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(1.9))
bullet(tf, "3DTK randomized 2-step Hough (baseline) · RANSAC primitive fit "
           "· nonlinear least-squares cylinder fit · Schnabel "
           "Efficient RANSAC (CGAL).", size=14, first=True, bold=True)
bullet(tf, "Architecture: a clustering proposer feeds a reusable metric fit; "
           "Efficient RANSAC is self-contained (no proposer) — isolating "
           "the proposer as a confound.", size=13, color=GREY)
bullet(tf, "The baseline emits a phantom secondary cylinder (panel 1 centre) "
           "and is non-deterministic; the seeded fitting methods stay clean.",
       size=13, color=GREY)
footer(s, 5)

# =====================================================================
# 6. RESULTS TABLE
# =====================================================================
s = slide()
title_band(s, "Quantitative accuracy on the real capture",
           "Identical metric applied to every method (axis-angle, radius "
           "error vs. known drum)")
rows = [
    ("Method", "Radius error", "Axis-angle error", "Runtime"),
    ("RANSAC primitive fit", "0.38 cm", "1.85°", "~0.7 s"),
    ("Nonlinear least-squares", "0.47 cm", "1.64°", "~9 s"),
    ("Efficient RANSAC (Schnabel)", "0.94 cm", "0.10°", "~0.05 s"),
    ("3DTK Hough (baseline)", "0.09 cm", "1.53°", "—"),
]
nrows, ncols = len(rows), len(rows[0])
gtbl = s.shapes.add_table(nrows, ncols, Inches(0.85), Inches(1.55),
                          Inches(11.6), Inches(3.0)).table
widths = [Inches(4.6), Inches(2.6), Inches(2.6), Inches(1.8)]
for i, w in enumerate(widths):
    gtbl.columns[i].width = w
for r in range(nrows):
    for c in range(ncols):
        cell = gtbl.cell(r, c)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = rows[r][c]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(15 if r == 0 else 14)
        run.font.bold = (r == 0)
        run.font.name = "Calibri"
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        else:
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
tf = box(s, Inches(0.85), Inches(4.9), Inches(11.6), Inches(2.1))
bullet(tf, "All methods recover radius to ~1 cm and axis to ~2°; "
           "operating points differ — Efficient RANSAC fastest, fitting "
           "methods most robust, the baseline over-reports (a phantom "
           "second cylinder).", size=14, first=True, bold=True, color=GREEN)
bullet(tf, "Single-barrel, low-clutter scene at ~180° visibility — the "
           "discriminating regime (occlusion, burial, clutter, multiple "
           "instances) is the next experiment.", size=13, italic=True,
       color=AMBER)
footer(s, 6)

# =====================================================================
# 7. STATUS
# =====================================================================
s = slide()
title_band(s, "Status: implemented vs. pending")
tf = box(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(5.4))
setp(tf.paragraphs[0], "Implemented", size=20, bold=True, color=GREEN)
bullet(tf, "Method-agnostic capture → detect → evaluate pipeline with "
           "a shared schema", size=14)
bullet(tf, "4 LiDAR-only geometric detectors (Hough, RANSAC, LS, "
           "Efficient RANSAC)", size=14)
bullet(tf, "Validated on synthetic ground truth + first real captures", size=14)
bullet(tf, "Literature survey of ~18 methods across the three families", size=14)

tf2 = box(s, Inches(6.9), Inches(1.45), Inches(6.0), Inches(5.4))
setp(tf2.paragraphs[0], "Pending", size=20, bold=True, color=AMBER)
bullet(tf2, "Occlusion sweep: recall and pose error vs. occlusion fraction "
            "(the headline experiment)", size=14)
bullet(tf2, "Learned detectors (PointNet++ / BarrelNet-style, occlusion-aware "
            "BtcDet) — require labelled training data", size=14)
bullet(tf2, "Multi-instance and cluttered real scenes (exercises the "
            "proposer)", size=14)
bullet(tf2, "Camera–LiDAR fusion pipeline (camera-first and LiDAR-first)",
       size=14)
bullet(tf2, "Validation on a real industrial test-site dataset", size=14,
       bold=True, color=ACCENT)
footer(s, 7)

# =====================================================================
# 8. SIMULATION + FUSION detail
# =====================================================================
s = slide()
title_band(s, "Simulation pipeline & fusion strategy")
tf = box(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(5.4))
setp(tf.paragraphs[0], "NVIDIA Isaac Sim — ground-truthed data engine",
     size=17, bold=True, color=ACCENT)
bullet(tf, "Replicator yields per-barrel pose/size/count, 3D boxes and "
           "instance/semantic masks for free.", size=14)
bullet(tf, "RTX LiDAR matched to the real sensor; perfectly co-registered "
           "camera (known extrinsics).", size=14)
bullet(tf, "Scripted scenes sweep barrel count, spacing, burial depth and "
           "mutual occlusion.", size=14)
bullet(tf, "True occlusion fraction (visible/total surface points) gives a "
           "grounded x-axis for the occlusion sweep.", size=14)
bullet(tf, "Supplies the labelled training set the learned detectors require.",
       size=14, color=GREEN)

tf2 = box(s, Inches(6.9), Inches(1.45), Inches(6.0), Inches(5.4))
setp(tf2.paragraphs[0], "Camera–LiDAR fusion", size=17, bold=True,
     color=ACCENT)
bullet(tf2, "Camera-first: image detection → frustum → constrained 3D "
            "fit (e.g. Frustum-PointNet).", size=14)
bullet(tf2, "LiDAR-first: 3D proposals → image verification / point "
            "painting.", size=14)
bullet(tf2, "Hypothesis to test, not assume: fusion's recall advantage under "
            "occlusion vs. LiDAR-only.", size=14)
bullet(tf2, "Calibration-noise injection: perturb the (sim-perfect) "
            "camera–LiDAR extrinsics to find where any fusion advantage "
            "degrades — a key real-world failure mode.", size=14,
       color=ACCENT)
footer(s, 8)

# =====================================================================
# 9. REQUIREMENT SPECIFICATIONS
# =====================================================================
s = slide()
title_band(s, "Requirement specifications & dataset request",
           "Inputs needed to parameterise the simulation and validate on "
           "real operating conditions")
tf = box(s, Inches(0.55), Inches(1.45), Inches(6.1), Inches(5.6))
setp(tf.paragraphs[0], "Target specification", size=16, bold=True,
     color=ACCENT)
bullet(tf, "What is the barrel size?  (nominal dimensions / drum standard, "
           "e.g. 200 L)", size=14, space_after=8)
bullet(tf, "Are the barrels damaged or discoloured?  (dents, rust, crushing; "
           "surface colour & finish)", size=14, space_after=8)
bullet(tf, "How much of each barrel is occluded?  (burial, stacking, clutter "
           "— i.e. typical visible fraction)", size=14, space_after=8)
add_para(tf, "Scene specification", size=16, bold=True, color=ACCENT,
         space_after=3)
bullet(tf, "How many barrels per scene, and their arrangement?", size=14)
bullet(tf, "Operating environment: indoor/outdoor, substrate, illumination, "
           "airborne dust.", size=14)

tf2 = box(s, Inches(6.9), Inches(1.45), Inches(6.0), Inches(5.6))
setp(tf2.paragraphs[0], "Sensor & data specification", size=16, bold=True,
     color=GREEN)
bullet(tf2, "Test-site recordings: LiDAR + camera — time-synced only if "
            "available; extrinsic calibration where available.", size=14,
       bold=True)
bullet(tf2, "Sensor suite: LiDAR model (channels / FoV / range), camera "
            "resolution; platform & mounting (robot / drone / fixed), motion.",
       size=14)
bullet(tf2, "Any annotations: barrel positions, counts, sizes (even "
            "approximate).", size=14)
add_para(tf2, "Operational requirements", size=16, bold=True, color=ACCENT,
         space_after=3)
bullet(tf2, "Required positional & radial accuracy, effective range, latency "
            "(online vs. offline), acceptable FP/FN trade-off.", size=13)
footer(s, 9)

# =====================================================================
# 10. SUMMARY
# =====================================================================
s = slide()
dark_bg(s)
band(s, ACCENT, Inches(0.18))
tf = box(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(4.8))
setp(tf.paragraphs[0], "Summary & next step", size=32, bold=True, color=WHITE)
add_para(tf, "A method-agnostic benchmark with four LiDAR-only geometric "
             "detectors is operational and quantified on synthetic and real "
             "data (radius ~1 cm, axis ~2°).", size=19,
         color=RGBColor(0xD6, 0xE6, 0xF0), space_after=12)
p = tf.add_paragraph()
setp(p, "• Next: occlusion/burial sweep via Isaac Sim, learned "
        "detectors, and camera–LiDAR fusion — all on the same metrics.",
     size=17, color=RGBColor(0xC9, 0xD6, 0xE0))
p = tf.add_paragraph()
setp(p, "• A calibrated test-site dataset would parameterise the "
        "simulation and validate under real operating conditions.", size=17,
     color=RGBColor(0xC9, 0xD6, 0xE0))

prs.save(OUT)
print("wrote", OUT)
print("slides:", len(prs.slides._sldIdLst))
