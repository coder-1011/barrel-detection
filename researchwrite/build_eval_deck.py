#!/usr/bin/env python3
"""
Build the methods-comparison EVALUATION deck (.pptx) — v2 after supervisor feedback.

Feedback applied (supervisor review of the 2026-07-01 draft):
  - much bigger fonts everywhere (body >= 18 pt, readable from the back row)
  - ONE point per bullet, plain language (no method jargon like "accumulator mass")
  - new general intro: the bin-picking problem, then "what if the objects are covered"
  - "fitting" explained with a line example + a cylinder-parameter figure
  - family pros/cons split into separate bullets, each family's idea explained roughly
  - important subtitle content promoted into highlighted key-lines
  - evaluation bullets integrated INTO the pipeline graphic (pipeline_flow.png)
  - synthetic experiment split into setup (with a picture of the data) + results
    (full-width plot); real data split into two slides (clean lab / occluded drum)
  - risks slide now pairs every risk with how we deal with it; positive summary at the end
  - NEW: three BarrelNet slides (method + synthetic training data, epoch comparison
    laptop-CPU vs A100, detection visualisation on the real pile)

v3 (2026-07-03, second feedback round from the user):
  - intro slide said the same thing 3x -> title + images + ONE keyline only
  - fitting explained BEFORE the covered-drums slide (general -> specific flow);
    the "small patch fits many cylinders" message now lands on the pile slide
  - NEW "test data" slide (synth / lab barrel / pile) BEFORE the methods slide;
    the synthetic-experiment setup slide is merged into the results slide
  - methods slide: architecture diagram + "details on request" subtitle removed
  - pipeline_flow.png regenerated without the italic per-box footnotes
  - later slides de-duplicated (BarrelNet caption vs bullet, limits vs Exp-2 noise
    bullet)

Run AFTER:
  - researchwrite/make_eval_figures.py   (intro/fitting/pipeline/synth-data/barrelnet figs)
  - eval/plot_noise_sweep.py             (synth_noise_sweep.png, big-font version)

Slide "synthetic results" bullets embed the F1 numbers read live from
eval/<method>.csv (sweep_* rows), so they are measured values, never hand-typed.

If the supervisor's AI-generated bin-picking images are saved as
presentation_assets/binpicking_ai_1.png + binpicking_ai_2.png they are used on the
intro slide automatically; otherwise the schematic binpicking_intro.png is used.

Output: researchwrite/evaluation_presentation.pptx
"""
import csv
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
MASTERS = os.path.abspath(os.path.join(HERE, ".."))
ASSETS = os.path.join(HERE, "presentation_assets")
BILF = os.path.join(HERE, "bilfinger_slides", "assets")
BNFIG = os.path.join(MASTERS, "methods", "barrelnet", "figures")
OUT = os.path.join(HERE, "evaluation_presentation.pptx")

# ---- title-slide metadata (brief: don't guess — left as placeholders) -------
THESIS_TITLE = "Detection and Robotic Manipulation of Partially Occluded Object(s)"
PROBLEM_LINE = ("Detecting partially occluded cylindrical objects for robotic "
                "manipulation: a multi-method comparison in 3D point clouds")
PRESENTER = "[FILL: presenter]"
SUPERVISOR = "[FILL: supervisor]"
DATELINE = "[FILL: date]"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

DARK = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x0B, 0x5C, 0x8A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
RED = RGBColor(0xB0, 0x2A, 0x2A)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
PALE = RGBColor(0xE3, 0xEE, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ----------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def setp(p, text, size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         space_after=8, level=0, italic=False):
    p.text = text
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"


def add_para(tf, *a, **k):
    p = tf.add_paragraph()
    setp(p, *a, **k)
    return p


def bullet(tf, text, size=20, color=DARK, level=0, bold=False, first=False,
           space_after=10, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    pre = "• " if level == 0 else "– "
    setp(p, pre + text, size=size, bold=bold, color=color, level=level,
         space_after=space_after, italic=italic)
    return p


def band(s, color=ACCENT, h=Inches(1.15)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def title_band(s, title, sub=None):
    band(s)
    tf = box(s, Inches(0.5), Inches(0.10), Inches(12.3), Inches(1.0))
    setp(tf.paragraphs[0], title, size=30, bold=True, color=WHITE)
    if sub:
        add_para(tf, sub, size=17, color=RGBColor(0xD6, 0xE6, 0xF0))


def keyline(s, text, t=Inches(6.55), color=ACCENT, size=20):
    """A highlighted one-line takeaway strip — the slide's key message."""
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), t,
                            Inches(12.55), Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb = PALE
    sh.line.color.rgb = color; sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    setp(tf.paragraphs[0], text, size=size, bold=True, color=color,
         align=PP_ALIGN.CENTER)


def footer(s, n):
    tf = box(s, Inches(0.4), Inches(7.08), Inches(12.0), Inches(0.35))
    setp(tf.paragraphs[0], THESIS_TITLE, size=10, color=GREY)
    tn = box(s, Inches(12.6), Inches(7.08), Inches(0.6), Inches(0.35))
    setp(tn.paragraphs[0], str(n), size=11, color=GREY, align=PP_ALIGN.RIGHT)


def pic_fit(s, path, l, t, w, h):
    from PIL import Image
    iw, ih = Image.open(path).size
    boxr, imgr = w / h, iw / ih
    if imgr > boxr:
        nw, nh = w, int(w / imgr)
    else:
        nh, nw = h, int(h * imgr)
    nl = l + (w - nw) // 2
    nt = t + (h - nh) // 2
    s.shapes.add_picture(path, Emu(int(nl)), Emu(int(nt)),
                         width=Emu(int(nw)), height=Emu(int(nh)))


def caption(s, text, l, t, w, color=GREY, size=14):
    tf = box(s, l, t, w, Inches(0.4))
    setp(tf.paragraphs[0], text, size=size, color=color, align=PP_ALIGN.CENTER)


def dark_bg(s):
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK
    bg.line.fill.background(); bg.shadow.inherit = False


def table(s, rows, l, t, w, h, col_widths, head_color=ACCENT, fontsize=15,
          head_fontsize=16):
    nr, nc = len(rows), len(rows[0])
    tbl = s.shapes.add_table(nr, nc, l, t, w, h).table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    for r in range(nr):
        for c in range(nc):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            p = cell.text_frame.paragraphs[0]
            p.text = str(rows[r][c])
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(head_fontsize if r == 0 else fontsize)
            run.font.bold = (r == 0)
            run.font.name = "Calibri"
            if r == 0:
                run.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = head_color
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
    return tbl


# ---------------- synthetic-sweep numbers, read live from eval/*.csv ---------
def _f(x):
    try:
        return float(x)
    except (TypeError, ValueError):
        return None


def sweep_f1(method):
    """Aggregate F1 over all sweep_* scenes of one method (live from CSV)."""
    path = os.path.join(MASTERS, "eval", f"{method}.csv")
    tp = fp = fn = 0
    if not os.path.isfile(path):
        return None
    with open(path) as f:
        for row in csv.DictReader(f):
            if row.get("scene", "").startswith("sweep_"):
                tp += int(_f(row["tp"]) or 0)
                fp += int(_f(row["fp"]) or 0)
                fn += int(_f(row["fn"]) or 0)
    if tp + fp + fn == 0:
        return None
    P = tp / (tp + fp) if (tp + fp) else 0.0
    R = tp / (tp + fn) if (tp + fn) else 0.0
    return 2 * P * R / (P + R) if (P + R) else 0.0


F1_LS = sweep_f1("ls_cylinder")
F1_RS = sweep_f1("ransac_cylinder")
F1_ER = sweep_f1("efficient_ransac")
F1_HG = sweep_f1("3dtk_hough")


def fmt_f1(v):
    return f"{v:.2f}" if v is not None else "n/a"


# =====================================================================
# TITLE
# =====================================================================
s = slide()
dark_bg(s)
band(s, ACCENT, Inches(0.18))
tf = box(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(2.6))
setp(tf.paragraphs[0], THESIS_TITLE, size=36, bold=True, color=WHITE)
add_para(tf, PROBLEM_LINE, size=21, color=RGBColor(0x9F, 0xC5, 0xDD),
         space_after=18)
tf2 = box(s, Inches(0.9), Inches(5.4), Inches(11.6), Inches(1.6))
setp(tf2.paragraphs[0], "Master's thesis  ·  evaluation review", size=19,
     color=RGBColor(0xC9, 0xD6, 0xE0))
add_para(tf2, f"{PRESENTER}   ·   supervisor: {SUPERVISOR}   ·   {DATELINE}",
         size=16, color=RGBColor(0xC9, 0xD6, 0xE0))

# =====================================================================
# 1. INTRO — THE BIN-PICKING PROBLEM
# =====================================================================
s = slide()
title_band(s, "The bin-picking problem")
ai1 = os.path.join(ASSETS, "binpicking_ai_1.png")
ai2 = os.path.join(ASSETS, "binpicking_ai_2.png")
if os.path.isfile(ai1) and os.path.isfile(ai2):
    pic_fit(s, ai1, Inches(0.4), Inches(1.45), Inches(6.2), Inches(4.9))
    pic_fit(s, ai2, Inches(6.75), Inches(1.45), Inches(6.2), Inches(4.9))
else:
    pic_fit(s, os.path.join(ASSETS, "binpicking_intro.png"),
            Inches(0.7), Inches(1.45), Inches(11.9), Inches(4.9))
keyline(s, "Before the robot can grasp anything, it must know exactly "
           "where each object is and how it lies.", t=Inches(6.45))
footer(s, 1)

# =====================================================================
# 2. WHAT "FITTING" MEANS (general, before our specific case)
# =====================================================================
s = slide()
title_band(s, "Our task: fit a cylinder to the visible points",
           "“Fitting” = choosing model parameters so the model matches "
           "the measured points as well as possible")
pic_fit(s, os.path.join(ASSETS, "fitting_explained.png"),
        Inches(0.4), Inches(1.55), Inches(12.55), Inches(4.7))
keyline(s, "With a full view, fitting is easy — the difficulty starts when "
           "most of the object is hidden.", t=Inches(6.45))
footer(s, 2)

# =====================================================================
# 3. OUR CASE — THE DRUMS ARE COVERED
# =====================================================================
s = slide()
title_band(s, "Our case is harder: the drums are covered",
           "Real scan: 200-litre steel drums, tumbled and partially buried")
pic_fit(s, os.path.join(ASSETS, "station1_pile_raw.png"),
        Inches(0.4), Inches(1.5), Inches(9.0), Inches(4.8))
tf = box(s, Inches(9.55), Inches(1.7), Inches(3.55), Inches(4.6))
bullet(tf, "Drums lie in arbitrary directions.", size=19, first=True)
bullet(tf, "Sand and other drums hide most of each drum.", size=19)
bullet(tf, "The sensor sees the pile from one side only.", size=19)
keyline(s, "Each drum shows only a small patch — and a small patch fits "
           "many different cylinders.", t=Inches(6.45), color=AMBER)
footer(s, 3)

# =====================================================================
# 4. EXISTING SOLUTIONS — THREE FAMILIES (plain language)
# =====================================================================
s = slide()
title_band(s, "Existing solutions — three families of methods",
           "From an 18-method literature survey; judged by how they cope with "
           "hidden objects")
fam = [
    ("Geometry-based", ACCENT,
     "Search the 3D points directly for cylinder shapes.",
     "Works out of the box — no training data needed.",
     "Needs enough visible points; touching drums get merged."),
    ("Learning-based", GREEN,
     "A neural network learns what partial barrels look like from many "
     "examples.",
     "Can recognise a barrel from a small fragment.",
     "Needs lots of labelled training examples first."),
    ("Camera + LiDAR", AMBER,
     "Combine colour images with the 3D points.",
     "If one sensor is blocked, the other still sees the object.",
     "Both sensors must be aligned precisely — errors break the pairing."),
]
x0, w = Inches(0.4), Inches(4.11)
for i, (name, col, idea, pro, con) in enumerate(fam):
    l = Emu(int(x0 + i * (w + Inches(0.11))))
    hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.5), w,
                            Inches(0.6))
    hd.fill.solid(); hd.fill.fore_color.rgb = col
    hd.line.fill.background(); hd.shadow.inherit = False
    setp(hd.text_frame.paragraphs[0], name, size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
    body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(2.1), w,
                              Inches(4.15))
    body.fill.solid(); body.fill.fore_color.rgb = LIGHT
    body.line.color.rgb = col; body.line.width = Pt(1.25)
    body.shadow.inherit = False
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = Pt(8)
    setp(tf.paragraphs[0], idea, size=17, color=DARK, space_after=14)
    add_para(tf, "+ " + pro, size=17, bold=True, color=GREEN, space_after=12)
    add_para(tf, "− " + con, size=17, bold=True, color=RED)
keyline(s, "We compare all three families on the same data — geometry "
           "first, since it needs no training data.", t=Inches(6.45))
footer(s, 4)

# =====================================================================
# 5. THE TEST DATA (before the methods — easy to hard)
# =====================================================================
s = slide()
title_band(s, "The data we test every method on",
           "Three data sets — difficulty increases left to right")
pic_fit(s, os.path.join(ASSETS, "test_data_overview.png"),
        Inches(0.3), Inches(1.6), Inches(12.7), Inches(4.6))
keyline(s, "Synthetic clouds tell us the best case; the real pile tells us "
           "the truth.", t=Inches(6.45))
footer(s, 5)

# =====================================================================
# 6. METHODS CHOSEN (4 geometric detectors)
# =====================================================================
s = slide()
title_band(s, "Four geometry-based methods implemented")
tf = box(s, Inches(0.5), Inches(1.8), Inches(12.4), Inches(4.4))
bullet(tf, "3DTK Hough (baseline) — every point votes for the cylinders it "
           "could lie on; the strongest vote wins.", size=21, space_after=22,
       first=True)
bullet(tf, "RANSAC fit — try many random small point samples, keep the "
           "cylinder most points agree with.", size=21, space_after=22)
bullet(tf, "Least-squares fit — start from a rough guess, then adjust the "
           "cylinder until it matches the points best.", size=21,
       space_after=22)
bullet(tf, "Efficient RANSAC — a faster RANSAC variant that searches the "
           "whole scene for shapes in one pass.", size=21, space_after=22)
keyline(s, "All four: LiDAR points only · no training data · identical input "
           "and output — directly comparable.", t=Inches(6.5))
footer(s, 6)

# =====================================================================
# 7. HOW WE EVALUATE
# =====================================================================
s = slide()
title_band(s, "How we evaluate — one pipeline for every method")
pic_fit(s, os.path.join(ASSETS, "pipeline_flow.png"),
        Inches(0.3), Inches(1.7), Inches(12.7), Inches(4.4))
keyline(s, "A prediction counts as correct if direction is within 30° and "
           "the centre within 10 cm of the true drum.", t=Inches(6.4))
footer(s, 7)

# =====================================================================
# 8. SYNTHETIC EXPERIMENT (setup folded into subtitle; data set shown
#    already on the test-data slide)
# =====================================================================
s = slide()
title_band(s, "Experiment 1 — how much noise can each method take?",
           "The synthetic barrel (120° visible) at 11 noise levels × "
           "3 repeats = 33 test clouds")
pic_fit(s, os.path.join(ASSETS, "synth_noise_sweep.png"),
        Inches(0.3), Inches(1.45), Inches(12.75), Inches(4.1))
tf = box(s, Inches(0.6), Inches(5.55), Inches(12.2), Inches(1.5))
bullet(tf, f"RANSAC and least-squares find the barrel at every noise level "
           f"(F1 = {fmt_f1(F1_RS)} / {fmt_f1(F1_LS)}).", size=18, first=True,
       color=GREEN, bold=True)
bullet(tf, "Efficient RANSAC degrades first; the Hough baseline is erratic — "
           "it even misses some clean scenes.", size=18, color=AMBER)
bullet(tf, "Least-squares is the most precise but ~5–10× slower than the "
           "others (~17 s per scene).", size=18, color=GREY)
footer(s, 8)

# =====================================================================
# 9. REAL DATA — CLEAN LAB BARREL
# =====================================================================
s = slide()
title_band(s, "Experiment 2 — the real lab barrel, no occlusion",
           "Depth-camera scan; true radius 4.25 cm")
pic_fit(s, os.path.join(BILF, "real_3d_det.png"),
        Inches(0.5), Inches(1.5), Inches(3.6), Inches(4.8))
caption(s, "Camera points (coloured) with the fitted cylinder (red)",
        Inches(0.2), Inches(6.35), Inches(4.4))
rows = [("Method", "Barrel found?", "Radius error", "Direction error"),
        ("RANSAC fit", "yes", "0.38 cm", "1.9°"),
        ("Least-squares", "yes", "0.47 cm", "1.6°"),
        ("Efficient RANSAC", "yes", "0.94 cm", "0.1°"),
        ("3DTK Hough", "yes + 1 false alarm", "0.09 cm", "1.5°")]
table(s, rows, Inches(4.6), Inches(1.75), Inches(8.35), Inches(2.9),
      [Inches(2.6), Inches(2.35), Inches(1.7), Inches(1.7)])
tf = box(s, Inches(4.7), Inches(4.9), Inches(8.2), Inches(1.4))
bullet(tf, "All four methods find the barrel.", size=19, first=True,
       color=GREEN, bold=True)
bullet(tf, "Errors are ~6–8× larger than on synthetic data — real sensor "
           "noise is not the clean noise we simulate.", size=18, color=AMBER)
footer(s, 9)

# =====================================================================
# 10. REAL DATA — OCCLUDED SURVEY DRUM
# =====================================================================
s = slide()
title_band(s, "Experiment 3 — a real occluded drum from the pile",
           "One tilted, half-buried 200 L drum cut out of the survey scan")
pic_fit(s, os.path.join(ASSETS, "real_fits_on_cloud.png"),
        Inches(0.35), Inches(1.5), Inches(7.3), Inches(4.5))
caption(s, "True drum (green) vs fitted cylinders on the real points",
        Inches(0.35), Inches(6.02), Inches(7.3))
tf = box(s, Inches(7.9), Inches(1.65), Inches(5.1), Inches(4.6))
bullet(tf, "All four methods fail the accuracy gate here.", size=19,
       first=True, bold=True, color=RED)
bullet(tf, "Best attempt (least-squares): direction good (17°), centre off "
           "by 11.1 cm — misses the 10 cm gate by 1.1 cm.", size=18)
bullet(tf, "The others: direction far off, or no detection at all.", size=18)
keyline(s, "Occlusion is exactly where the classical methods break — "
           "motivating a learned method.", t=Inches(6.48), color=AMBER)
footer(s, 10)

# =====================================================================
# 11. BARRELNET — THE FIRST LEARNED METHOD
# =====================================================================
s = slide()
title_band(s, "New: BarrelNet — our first learning-based method",
           "A neural network that gets the points of ONE drum and outputs its "
           "pose (centre + direction)")
pic_fit(s, os.path.join(BNFIG, "synth_patches_sample.png"),
        Inches(0.4), Inches(1.6), Inches(12.55), Inches(3.6))
caption(s, "Examples of the synthetic training patches: simulated LiDAR "
           "points (red) on the drum surface (blue)",
        Inches(0.4), Inches(5.15), Inches(12.55))
tf = box(s, Inches(0.6), Inches(5.55), Inches(12.2), Inches(0.9))
bullet(tf, "Trained on 12,000 such patches — buried, tilted, noisy, "
           "incomplete.", size=19, first=True)
keyline(s, "No real data in training: the 21 hand-verified real drums are "
           "kept aside as an honest test.", t=Inches(6.5), color=GREEN)
footer(s, 11)

# =====================================================================
# 12. BARRELNET — TRAINING PROGRESS (epoch comparison)
# =====================================================================
s = slide()
title_band(s, "BarrelNet — how the score develops during training",
           "Same network trained twice: laptop CPU (stopped early) vs GPU "
           "server (full schedule)")
pic_fit(s, os.path.join(ASSETS, "barrelnet_epochs.png"),
        Inches(0.3), Inches(1.5), Inches(12.75), Inches(4.15))
tf = box(s, Inches(0.6), Inches(5.6), Inches(12.2), Inches(1.4))
bullet(tf, "Finishing the full training schedule nearly doubled the score: "
           "7 → 12 of 21 real drums.", size=19, first=True, bold=True,
       color=GREEN)
bullet(tf, "The direction is learned early; the centre position improves "
           "late — it crosses the 10 cm gate only near the end.", size=18)
footer(s, 12)

# =====================================================================
# 13. BARRELNET — RESULT ON THE REAL PILE
# =====================================================================
s = slide()
title_band(s, "BarrelNet on the real drum pile",
           "Trained purely on synthetic drums, tested on the 21 real ones")
pic_fit(s, os.path.join(BNFIG, "station_detection.png"),
        Inches(0.3), Inches(1.45), Inches(8.2), Inches(5.0))
caption(s, "Annotated drums coloured; predicted cylinders: green = within "
           "gate, red = miss", Inches(0.3), Inches(6.5), Inches(8.2))
tf = box(s, Inches(8.75), Inches(1.65), Inches(4.3), Inches(5.0))
bullet(tf, "12 of 21 real drums located within the accuracy gate.", size=19,
       first=True, bold=True, color=GREEN)
bullet(tf, "Direction right on 16 of 21 drums.", size=19)
bullet(tf, "The geometric methods scored 0 on this kind of drum.", size=19)
bullet(tf, "Open point: BarrelNet estimates the pose of a drum it is given — "
           "finding the drums in the pile still needs a separate first step.",
       size=17, color=GREY)
footer(s, 13)

# =====================================================================
# 14. WHAT LIMITS US TODAY
# =====================================================================
s = slide()
title_band(s, "What limits the evaluation today")
tf = box(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(4.7))
bullet(tf, "Our synthetic scenes still cannot imitate a truly buried, "
           "cluttered pile.", size=21, first=True)
bullet(tf, "Real sensor noise is harder than anything we simulate — "
           "Experiment 2 showed 6–8× larger errors.", size=21)
bullet(tf, "Only 21 real drums are labelled (~30% of the pile) — too few for "
           "strong statistics.", size=21)
bullet(tf, "Some labels are imperfect: two drums lying end-to-end were "
           "annotated as one.", size=21)
keyline(s, "The evaluation machinery works — what is missing is more and "
           "better ground-truth data.", t=Inches(6.45))
footer(s, 14)

# =====================================================================
# 15. RISKS — AND HOW WE ADDRESS THEM (before the roadmap, so the deck
#     ends problems -> solutions -> summary)
# =====================================================================
s = slide()
title_band(s, "Open risks — and how we deal with them")
rows = [
    ("Risk", "How we deal with it"),
    ("Accuracy collapses under occlusion",
     "Measure it per method — that comparison IS the thesis."),
    ("One sensor alone may not be enough",
     "Add camera + LiDAR fusion — it degrades far more gracefully."),
    ("Sim-trained models may not transfer",
     "Measured: 12 of 21 real drums; next, fine-tune on a few real ones."),
    ("Manipulation not started yet",
     "Detection delivers exactly the pose a grasp planner needs."),
]
table(s, rows, Inches(0.45), Inches(1.6), Inches(12.45), Inches(4.2),
      [Inches(5.2), Inches(7.25)], fontsize=18, head_fontsize=19)
caption(s, "Occlusion numbers: [Shi'19] arXiv:1812.04244 · [Kumar'25] "
           "arXiv:2511.04347 · [Wang'25] Sensors 25(9):2794",
        Inches(0.45), Inches(6.45), Inches(12.45), size=12)
footer(s, 15)

# =====================================================================
# 16. FUTURE PIPELINE — SIMULATION CLOSES THE DATA GAP
# =====================================================================
s = slide()
title_band(s, "Next step: simulation closes the data gap",
           "NVIDIA Isaac Sim gives unlimited scenes with perfect ground truth")
pic_fit(s, os.path.join(ASSETS, "roadmap.png"),
        Inches(0.4), Inches(1.6), Inches(12.5), Inches(2.6))
tf = box(s, Inches(0.6), Inches(4.35), Inches(12.2), Inches(2.0))
bullet(tf, "Every simulated drum comes with its exact pose and its exact "
           "degree of burial — labels for free.", size=20, first=True)
bullet(tf, "That enables the headline experiment: detection quality as a "
           "function of how much of the drum is hidden.", size=20)
bullet(tf, "And it unlocks the remaining families: more learned methods and "
           "camera + LiDAR fusion.", size=20)
footer(s, 16)

# =====================================================================
# 17. SUMMARY
# =====================================================================
s = slide()
title_band(s, "Summary")
tf = box(s, Inches(0.55), Inches(1.6), Inches(12.3), Inches(4.6))
bullet(tf, "One fair evaluation pipeline — every method, same data, same "
           "score.", size=22, first=True, bold=True)
bullet(tf, "Four geometry-based methods work well on visible barrels — and "
           "all fail on a truly occluded one.", size=22)
bullet(tf, "BarrelNet, trained only on synthetic drums, already finds 12 of "
           "21 real occluded drums.", size=22, color=GREEN)
bullet(tf, "Next: simulation-generated data for the occlusion sweep, "
           "fine-tuning, and sensor fusion.", size=22)
keyline(s, "Learned + simulated is the promising path for occluded drums — "
           "and we can now measure exactly how promising.", t=Inches(6.45),
        color=GREEN)
footer(s, 17)

# =====================================================================
# 18. THANK YOU
# =====================================================================
s = slide()
dark_bg(s)
band(s, ACCENT, Inches(0.18))
tf = box(s, Inches(0.9), Inches(2.8), Inches(11.6), Inches(2.0))
setp(tf.paragraphs[0], "Thank you!", size=48, bold=True, color=WHITE,
     align=PP_ALIGN.CENTER)
add_para(tf, "Questions?", size=24, color=RGBColor(0x9F, 0xC5, 0xDD),
         align=PP_ALIGN.CENTER)

prs.save(OUT)
print("wrote", OUT)
print("slides:", len(prs.slides._sldIdLst))
